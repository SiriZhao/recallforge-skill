from __future__ import annotations

import re
from pathlib import Path

from .classify import classify_document
from .config import SUPPORTED_EXTENSIONS
from .models import Document, DocumentBlock, GenerationReport
from .ocr import ocr_image


def _split_text_blocks(source_file: str, text: str, doc_type: str) -> list[DocumentBlock]:
    blocks: list[DocumentBlock] = []
    current_heading = None
    current_chapter = None
    page = None
    for raw in re.split(r"\n\s*\n", text.replace("\r\n", "\n")):
        content = raw.strip()
        if not content:
            continue
        heading_match = re.match(r"^(第[一二三四五六七八九十\d]+章|#+\s*.+|[0-9]+[.、]\s*.+)", content)
        if heading_match:
            current_heading = content.splitlines()[0].lstrip("# ").strip()
        chap_match = re.search(r"第[一二三四五六七八九十\d]+章\s*([^\n：:]+)?", content)
        if chap_match:
            current_chapter = chap_match.group(0)
        page_match = re.search(r"(?:page|页|slide|幻灯片)\s*[:：]?\s*(\d+)", content, re.I)
        if page_match:
            page = page_match.group(1)
        q_match = re.search(r"(?:第\s*)?(\d+)\s*[题、.]", content)
        blocks.append(DocumentBlock(
            source_file=source_file,
            content=content,
            page_or_slide=page,
            question_number=q_match.group(1) if q_match else None,
            doc_type=doc_type,
            chapter=current_chapter,
            heading=current_heading,
            source_refs=[{"source_file": source_file, "page_or_slide": page, "question_number": q_match.group(1) if q_match else None, "heading": current_heading}],
        ))
    return blocks


def ingest_path(input_dir: Path, report: GenerationReport) -> list[Document]:
    docs: list[Document] = []
    files = sorted([p for p in input_dir.rglob("*") if p.is_file() and p.suffix.lower() in SUPPORTED_EXTENSIONS])
    report.files_seen = [str(p) for p in files]
    ocr_cache = input_dir / ".ocr_cache"
    for path in files:
        warnings: list[str] = []
        try:
            ext = path.suffix.lower()
            text = ""
            if ext in {".txt", ".md"}:
                text = path.read_text(encoding="utf-8", errors="ignore")
            elif ext == ".pdf":
                try:
                    from pypdf import PdfReader

                    reader = PdfReader(str(path))
                    parts = []
                    for i, page in enumerate(reader.pages, 1):
                        parts.append(f"Page {i}\n{page.extract_text() or ''}")
                    text = "\n\n".join(parts)
                except Exception as exc:
                    warnings.append(f"PDF parser unavailable or failed for {path.name}: {exc}. 若为扫描 PDF，请转图片或安装 OCR 依赖。")
            elif ext == ".docx":
                try:
                    import docx

                    d = docx.Document(str(path))
                    text = "\n\n".join(p.text for p in d.paragraphs if p.text.strip())
                except Exception as exc:
                    warnings.append(f"DOCX parser unavailable or failed for {path.name}: {exc}.")
            elif ext == ".pptx":
                try:
                    from pptx import Presentation

                    prs = Presentation(str(path))
                    parts = []
                    for i, slide in enumerate(prs.slides, 1):
                        lines = [f"Slide {i}"]
                        for shape in slide.shapes:
                            if hasattr(shape, "text") and shape.text.strip():
                                lines.append(shape.text)
                        parts.append("\n".join(lines))
                    text = "\n\n".join(parts)
                except Exception as exc:
                    warnings.append(f"PPTX parser unavailable or failed for {path.name}: {exc}.")
            elif ext in {".png", ".jpg", ".jpeg"}:
                text, conf = ocr_image(path, ocr_cache, warnings)
                if conf < 0.5:
                    warnings.append(f"OCR low confidence for {path.name}; 低置信度内容不会作为高置信度结论。")
            doc_type = classify_document(path, text)
            blocks = _split_text_blocks(path.name, text, doc_type)
            doc = Document(source_file=path.name, doc_type=doc_type, blocks=blocks, warnings=warnings)
            docs.append(doc)
            report.files_read.append(str(path))
            for w in warnings:
                report.warn(w)
        except Exception as exc:
            report.warn(f"Failed to ingest {path.name}: {exc}. 已跳过该文件。")
    if not files:
        report.warn(f"No supported files found in {input_dir}.")
    return docs
