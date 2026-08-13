from __future__ import annotations

import re
from pathlib import Path

from .types import NativePage


FORMULA_TOKEN = re.compile(
    r"[=≠≈≤≥]|[∑∫√∂Δπθαβγμλσφψω]|\\frac|\\sum|\\int|_[a-zA-Z0-9]|\^[a-zA-Z0-9]"
)
FRACTION_RE = re.compile(r"\d+\s*/\s*\d+")
SUBSCRIPT_RE = re.compile(r"(?<![\w])([A-Za-z]{1,3})[_ ]([0-9]+)")
SUPERSCRIPT_RE = re.compile(r"(?<![\w])([A-Za-z]{1,3})[\^] ?([0-9]+)")


def _language_hint(text: str) -> str | None:
    """Cheap first-pass language hint from a text sample (zh / en / mixed)."""
    sample = text[:2000]
    han = sum(1 for ch in sample if "\u4e00" <= ch <= "\u9fff")
    latin = sum(1 for ch in sample if ch.isascii() and ch.isalpha())
    if han > 0 and latin == 0:
        return "zh-CN"
    if latin > 0 and han == 0:
        return "en-US"
    if han > 0 and latin > 0:
        return "mixed"
    return None


def _parse_text_file(path: Path) -> list[NativePage]:
    text = path.read_text(encoding="utf-8", errors="ignore")
    page = NativePage(
        page_or_slide="1",
        raw_text=text,
        has_text_layer=True,
        heading=path.stem,
        language_hint=_language_hint(text),
    )
    return [page]


def _parse_pdf(path: Path) -> list[NativePage]:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    pages: list[NativePage] = []
    for i, pdf_page in enumerate(reader.pages, 1):
        try:
            text = pdf_page.extract_text() or ""
        except Exception:
            text = ""
        images = len(pdf_page.images or [])
        page = NativePage(
            page_or_slide=str(i),
            raw_text=text,
            has_text_layer=bool(text.strip()),
            has_images=images > 0,
            language_hint=_language_hint(text),
        )
        _detect_formula_signals(page)
        _detect_questions(page)
        pages.append(page)
    return pages


def _parse_pdf_pymupdf(path: Path) -> list[NativePage]:
    """PDF native parsing via PyMuPDF: richer text spans and image detection."""
    import pymupdf

    doc = pymupdf.open(str(path))
    pages: list[NativePage] = []
    for i, page in enumerate(doc, 1):
        text = page.get_text("text") or ""
        images = page.get_images(full=True) or []
        np = NativePage(
            page_or_slide=str(i),
            raw_text=text,
            has_text_layer=bool(text.strip()),
            has_images=bool(images),
            language_hint=_language_hint(text),
        )
        # super/subscript detection from span font sizes (formula-heavy signal)
        try:
            for block in page.get_text("dict").get("blocks", []):
                for line in block.get("lines", []):
                    sizes = [s.get("size", 0) for s in line.get("spans", [])]
                    if len(sizes) >= 2 and max(sizes) - min(sizes) >= 1.0:
                        if "superscript/subscript" not in np.formula_signals:
                            np.formula_signals.append("superscript/subscript")
        except Exception:
            pass
        _detect_formula_signals(np)
        _detect_questions(np)
        pages.append(np)
    doc.close()
    return pages


def _detect_formula_signals(page: NativePage) -> None:
    text = page.raw_text
    if FORMULA_TOKEN.search(text):
        page.formula_signals.append("math-tokens")
    if FRACTION_RE.search(text):
        page.formula_signals.append("fraction")
    if SUBSCRIPT_RE.search(text) and "superscript/subscript" not in page.formula_signals:
        page.formula_signals.append("subscript")
    if SUPERSCRIPT_RE.search(text) and "superscript/subscript" not in page.formula_signals:
        page.formula_signals.append("superscript")
    if re.search(r"\b(?:NaOH|HCl|H2O|CO2|NH3|KMnO4|CaCO3|CH4|C2H5OH)\b", text, re.I):
        page.formula_signals.append("chemical-equation")


def _detect_questions(page: NativePage) -> None:
    text = page.raw_text
    found = re.findall(r"(?:^|\n)\s*(?:第\s*)?(\d+)\s*[.、．]?\s*(?:题)?", text)
    page.question_numbers = list(dict.fromkeys(found))[:20]


def _parse_pptx(path: Path) -> list[NativePage]:
    from pptx import Presentation

    prs = Presentation(str(path))
    pages: list[NativePage] = []
    for i, slide in enumerate(prs.slides, 1):
        lines: list[str] = []
        has_images = False
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.text.strip():
                lines.append(shape.text)
            if getattr(shape, "shape_type", None) is not None and "PICTURE" in str(shape.shape_type):
                has_images = True
        text = "\n".join(lines)
        page = NativePage(
            page_or_slide=str(i),
            raw_text=text,
            has_text_layer=bool(text.strip()),
            has_images=has_images,
            heading=_slide_title(slide),
            language_hint=_language_hint(text),
        )
        _detect_formula_signals(page)
        _detect_questions(page)
        pages.append(page)
    return pages


def _slide_title(slide) -> str | None:
    try:
        for shape in slide.shapes:
            if shape.has_text_frame and shape == slide.shapes.title:
                return shape.text.strip() or None
    except Exception:
        pass
    return None


def _parse_docx(path: Path) -> list[NativePage]:
    import docx

    document = docx.Document(str(path))
    parts: list[str] = []
    for para in document.paragraphs:
        if para.text.strip():
            parts.append(para.text)
    for table in document.tables:
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            parts.append(" | ".join(cells))
    text = "\n".join(parts)
    page = NativePage(
        page_or_slide="1",
        raw_text=text,
        has_text_layer=bool(text.strip()),
        heading=path.stem,
        table_signals=bool(document.tables),
        language_hint=_language_hint(text),
    )
    _detect_formula_signals(page)
    _detect_questions(page)
    return [page]


def _parse_image(path: Path) -> list[NativePage]:
    # No native text for images: the visual path is required. No OCR by default.
    return [
        NativePage(
            page_or_slide="1",
            raw_text="",
            has_text_layer=False,
            has_images=True,
            heading=path.stem,
        )
    ]


def parse_native(path: Path, document_type: str, *, prefer_pymupdf: bool = True) -> list[NativePage]:
    """Native-first parsing. Raises on genuinely unreadable files (never fabricates)."""
    if document_type in {"txt", "md"}:
        return _parse_text_file(path)
    if document_type == "pdf":
        if prefer_pymupdf:
            try:
                return _parse_pdf_pymupdf(path)
            except Exception:
                return _parse_pdf(path)
        return _parse_pdf(path)
    if document_type == "pptx":
        return _parse_pptx(path)
    if document_type == "docx":
        return _parse_docx(path)
    if document_type in {"png", "jpg"}:
        return _parse_image(path)
    raise ValueError(f"no native parser for {document_type}")
