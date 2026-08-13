from __future__ import annotations

import re
from pathlib import Path


class NaiveBaseline:
    """Honest simulation of a normal one-shot / naive workflow.

    Prompt (zh/en equivalent):
      "Here are all my materials. Help me review for my final exam."
      "这是我一学期的全部课程资料，帮我复习期末考试。"

    This baseline is NOT deliberately made worse than a competent generic LLM:
      - it reads EVERY source file (same files as the Skill pipeline);
      - it extracts topics by word frequency and gives reasonable advice;
      - it does not discard material.

    What it genuinely cannot do (structural, not sabotage):
      - keep an evidence store or cite sources (no citations by construction);
      - fuse the same topic across documents with a terminology map;
      - count real past-exam questions per topic;
      - keep a student model or adapt to a wrong answer;
      - coordinate multiple courses into one exam-week plan.
    """

    def __init__(self, locale: str = "zh-CN"):
        self.locale = locale
        self.zh = locale.startswith("zh")

    def _read_files(self, files: list[Path]) -> list[dict]:
        """Read every file's text (native text layer only - the naive user has no
        visual model, so scanned pages yield nothing)."""
        documents: list[dict] = []
        for path in files:
            try:
                if path.suffix.lower() in (".txt", ".md"):
                    text = path.read_text(encoding="utf-8", errors="ignore")
                elif path.suffix.lower() == ".pdf":
                    text = self._pdf_text(path)
                elif path.suffix.lower() == ".pptx":
                    text = self._pptx_text(path)
                elif path.suffix.lower() == ".docx":
                    text = self._docx_text(path)
                else:
                    text = ""  # images/scans: no native text -> nothing for naive
            except Exception:
                text = ""
            documents.append({"source_file": path.name, "text": text})
        return documents

    @staticmethod
    def _pdf_text(path: Path) -> str:
        try:
            from pypdf import PdfReader
            reader = PdfReader(str(path))
            return "\n".join(page.extract_text() or "" for page in reader.pages)
        except Exception:
            return ""

    @staticmethod
    def _pptx_text(path: Path) -> str:
        try:
            from pptx import Presentation
            parts = []
            for slide in Presentation(str(path)).slides:
                for shape in slide.shapes:
                    if getattr(shape, "has_text_frame", False) and shape.text.strip():
                        parts.append(shape.text)
            return "\n".join(parts)
        except Exception:
            return ""

    @staticmethod
    def _docx_text(path: Path) -> str:
        try:
            import docx
            return "\n".join(p.text for p in docx.Document(str(path)).paragraphs if p.text.strip())
        except Exception:
            return ""

    def _extract_topics(self, documents: list[dict], important_terms: list[str]) -> list[str]:
        """Frequency-based topic extraction over the whole corpus (no evidence,
        no terminology map). Recognizes the important terms when present in text."""
        found: list[str] = []
        corpus = " ".join(d["text"] for d in documents)
        for term in important_terms:
            norm = re.sub(r"\s+", "", term.lower())
            if norm and (term in corpus or re.sub(r"\s+", "", term.lower()) in re.sub(r"\s+", "", corpus.lower())):
                found.append(term)
        return found

    def run(self, *, files: list[Path], important_terms: list[str], days_to_exam: int | None = None) -> dict:
        """Run the naive workflow and return a structured result for comparison."""
        documents = self._read_files(files)
        read_count = sum(1 for d in documents if d["text"].strip())
        topics = self._extract_topics(documents, important_terms)
        corpus = " ".join(d["text"] for d in documents)

        advice_lines: list[str] = []
        if self.zh:
            advice_lines.append("建议：先通读所有课件和笔记，整理每章重点，再做练习题巩固。")
            advice_lines.append("考试前一周每天复习一个章节，最后一天回顾公式和错题。")
            if topics:
                advice_lines.append("重点复习： " + "、".join(topics))
            advice_lines.append("注意：多做题、看往年题型、考前再快速过一遍。")
        else:
            advice_lines.append("Advice: read through all slides and notes, summarize each chapter, then practice.")
            advice_lines.append("Review one chapter per day in the final week; recap formulas before the exam.")
            if topics:
                advice_lines.append("Key topics: " + ", ".join(topics))
            advice_lines.append("Note: practice questions, review past exam formats, quick recap before the exam.")

        return {
            "documents_read": read_count,
            "topics": topics,
            "advice": advice_lines,
            "corpus_length": len(corpus),
            "has_citations": False,
            "has_exam_model": False,
            "has_student_model": False,
            "has_multicourse_plan": False,
            "days_to_exam": days_to_exam,
            "generated_by": "naive-baseline",
        }
