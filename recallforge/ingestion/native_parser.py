from __future__ import annotations

import hashlib
import re
from pathlib import Path

from .types import NativePage


FORMULA_TOKEN = re.compile(
    r"[=≠≈≤≥]|[∑∫√∂Δπθαβγμλσφψω]|\\frac|\\sum|\\int|_[a-zA-Z0-9]|\^[a-zA-Z0-9]"
)
FRACTION_RE = re.compile(r"\d+\s*/\s*\d+")
SUBSCRIPT_RE = re.compile(r"(?<![\w])([A-Za-z]{1,3})[_ ]([0-9]+)")
SUPERSCRIPT_RE = re.compile(r"(?<![\w])([A-Za-z]{1,3})[\^] ?([0-9]+)")


def _language_hint(text: str) -> str | None:
    """Cheap first-pass language hint from a text sample (zh / en / mixed)."""
    sample = text[:2000]
    han = sum(1 for ch in sample if "\u4e00" <= ch <= "\u9fff")
    latin = sum(1 for ch in sample if ch.isascii() and ch.isalpha())
    if han > 0 and latin == 0:
        return "zh-CN"
    if latin > 0 and han == 0:
        return "en-US"
    if han > 0 and latin > 0:
        return "mixed"
    return None


def _parse_text_file(path: Path) -> list[NativePage]:
    text = path.read_text(encoding="utf-8", errors="ignore")
    page = NativePage(
        page_or_slide="1",
        raw_text=text,
        has_text_layer=True,
        heading=path.stem,
        language_hint=_language_hint(text),
        blocks=[{"type": "text", "text": text, "order": 0}],
        native_confidence=0.98,
        source_anchor=f"{path.name}, p. 1",
        page_hash=hashlib.sha256(text.encode("utf-8")).hexdigest(),
    )
    return [page]


def _parse_pdf(path: Path) -> list[NativePage]:
    from pypdf import PdfReader

    reader = PdfReader(str(path))
    pages: list[NativePage] = []
    for i, pdf_page in enumerate(reader.pages, 1):
        try:
            text = pdf_page.extract_text() or ""
        except Exception:
            text = ""
        images = len(pdf_page.images or [])
        page = NativePage(
            page_or_slide=str(i),
            raw_text=text,
            has_text_layer=bool(text.strip()),
            has_images=images > 0,
            language_hint=_language_hint(text),
        )
        _detect_formula_signals(page)
        _detect_questions(page)
        pages.append(page)
    return pages


def _parse_pdf_pymupdf(path: Path) -> list[NativePage]:
    """PDF native parsing via PyMuPDF: richer text spans and image detection."""
    import pymupdf

    doc = pymupdf.open(str(path))
    pages: list[NativePage] = []
    for i, page in enumerate(doc, 1):
        text = page.get_text("text") or ""
        images = page.get_images(full=True) or []
        page_dict = page.get_text("dict")
        blocks: list[dict] = []
        for block in page_dict.get("blocks", []):
            if block.get("type") == 0:
                block_text = "\n".join(
                    "".join(span.get("text", "") for span in line.get("spans", []))
                    for line in block.get("lines", [])
                ).strip()
                if block_text:
                    blocks.append({"type": "text", "text": block_text, "bbox": block.get("bbox"), "order": len(blocks)})
            elif block.get("type") == 1:
                blocks.append({"type": "image", "bbox": block.get("bbox"), "order": len(blocks)})
        page_area = max(float(page.rect.width * page.rect.height), 1.0)
        image_area = sum(
            max(0.0, float(b.get("bbox", [0, 0, 0, 0])[2] - b.get("bbox", [0, 0, 0, 0])[0]))
            * max(0.0, float(b.get("bbox", [0, 0, 0, 0])[3] - b.get("bbox", [0, 0, 0, 0])[1]))
            for b in blocks if b.get("type") == "image"
        )
        printable = sum(1 for ch in text if ch.isprintable() and not ch.isspace())
        suspicious = sum(1 for ch in text if ch == "\ufffd" or (ord(ch) < 32 and ch not in "\n\r\t"))
        confidence = min(0.99, len(text.strip()) / 120.0) * (1.0 - min(0.9, suspicious / max(printable, 1)))
        np = NativePage(
            page_or_slide=str(i),
            raw_text=text,
            has_text_layer=bool(text.strip()),
            has_images=bool(images),
            language_hint=_language_hint(text),
            blocks=blocks,
            native_confidence=round(confidence, 3),
            image_coverage=round(min(1.0, image_area / page_area), 3),
            suspicious_char_ratio=round(suspicious / max(printable, 1), 3),
            rotation=int(page.rotation or 0),
            source_anchor=f"{path.name}, p. {i}",
            page_hash=hashlib.sha256((text + repr(blocks)).encode("utf-8")).hexdigest(),
        )
        # super/subscript detection from span font sizes (formula-heavy signal)
        try:
            for block in page.get_text("dict").get("blocks", []):
                for line in block.get("lines", []):
                    sizes = [s.get("size", 0) for s in line.get("spans", [])]
                    if len(sizes) >= 2 and max(sizes) - min(sizes) >= 1.0:
                        if "superscript/subscript" not in np.formula_signals:
                            np.formula_signals.append("superscript/subscript")
        except Exception:
            pass
        _detect_formula_signals(np)
        _detect_questions(np)
        pages.append(np)
    # Repeated short top/bottom lines are document furniture, not course evidence.
    repeated: dict[str, int] = {}
    for candidate in pages:
        lines = [line.strip() for line in candidate.raw_text.splitlines() if line.strip()]
        for line in (lines[:2] + lines[-2:]):
            if len(line) <= 100:
                repeated[line] = repeated.get(line, 0) + 1
    threshold = max(3, int(len(pages) * 0.6))
    repeated_furniture = {line for line, count in repeated.items() if count >= threshold}
    if repeated_furniture:
        for candidate in pages:
            candidate.raw_text = "\n".join(
                line for line in candidate.raw_text.splitlines() if line.strip() not in repeated_furniture
            ).strip()
            candidate.blocks = [
                block for block in candidate.blocks
                if block.get("type") != "text" or block.get("text", "").strip() not in repeated_furniture
            ]
            candidate.warning = "repeated header/footer removed"
    doc.close()
    return pages


def _detect_formula_signals(page: NativePage) -> None:
    text = page.raw_text
    if FORMULA_TOKEN.search(text):
        page.formula_signals.append("math-tokens")
    if FRACTION_RE.search(text):
        page.formula_signals.append("fraction")
    if SUBSCRIPT_RE.search(text) and "superscript/subscript" not in page.formula_signals:
        page.formula_signals.append("subscript")
    if SUPERSCRIPT_RE.search(text) and "superscript/subscript" not in page.formula_signals:
        page.formula_signals.append("superscript")
    if re.search(r"\b(?:NaOH|HCl|H2O|CO2|NH3|KMnO4|CaCO3|CH4|C2H5OH)\b", text, re.I):
        page.formula_signals.append("chemical-equation")


def _detect_questions(page: NativePage) -> None:
    text = page.raw_text
    found = re.findall(r"(?:^|\n)\s*(?:第\s*)?(\d+)\s*[.、．]?\s*(?:题)?", text)
    page.question_numbers = list(dict.fromkeys(found))[:20]


def _parse_pptx(path: Path) -> list[NativePage]:
    from pptx import Presentation

    prs = Presentation(str(path))
    pages: list[NativePage] = []
    for i, slide in enumerate(prs.slides, 1):
        lines: list[str] = []
        has_images = False
        blocks: list[dict] = []
        notes: list[str] = []
        emphasis: list[dict] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.text.strip():
                lines.append(shape.text)
                blocks.append({
                    "type": "text",
                    "text": shape.text,
                    "bbox": [int(shape.left), int(shape.top), int(shape.left + shape.width), int(shape.top + shape.height)],
                    "order": len(blocks),
                    "shape_type": str(getattr(shape, "shape_type", "")),
                })
                lowered = shape.text.lower()
                if any(marker in lowered for marker in ("important", "must know", "exam", "重点", "考试要求")):
                    emphasis.append({"kind": "explicit_text", "text": shape.text[:200]})
            if getattr(shape, "has_table", False):
                rows = [[cell.text for cell in row.cells] for row in shape.table.rows]
                blocks.append({"type": "table", "rows": rows, "headers": rows[0] if rows else [], "order": len(blocks)})
            if getattr(shape, "has_chart", False):
                blocks.append({"type": "chart", "order": len(blocks), "bbox": [int(shape.left), int(shape.top), int(shape.left + shape.width), int(shape.top + shape.height)]})
            if getattr(shape, "shape_type", None) is not None and "PICTURE" in str(shape.shape_type):
                has_images = True
                blocks.append({"type": "image", "order": len(blocks), "bbox": [int(shape.left), int(shape.top), int(shape.left + shape.width), int(shape.top + shape.height)]})
        try:
            note_text = slide.notes_slide.notes_text_frame.text.strip()
            if note_text:
                notes.append(note_text)
        except Exception:
            pass
        text = "\n".join(lines)
        page = NativePage(
            page_or_slide=str(i),
            raw_text=text,
            has_text_layer=bool(text.strip()),
            has_images=has_images,
            heading=_slide_title(slide),
            language_hint=_language_hint(text),
            blocks=blocks,
            notes=notes,
            native_confidence=0.92 if text.strip() else 0.0,
            source_anchor=f"{path.name}, slide {i}",
            visual_emphasis=emphasis,
            page_hash=hashlib.sha256((text + repr(blocks) + repr(notes)).encode("utf-8")).hexdigest(),
        )
        _detect_formula_signals(page)
        _detect_questions(page)
        pages.append(page)
    return pages


def _slide_title(slide) -> str | None:
    try:
        for shape in slide.shapes:
            if shape.has_text_frame and shape == slide.shapes.title:
                return shape.text.strip() or None
    except Exception:
        pass
    return None


def _parse_docx(path: Path) -> list[NativePage]:
    import docx

    document = docx.Document(str(path))
    parts: list[str] = []
    for para in document.paragraphs:
        if para.text.strip():
            parts.append(para.text)
    blocks = [{"type": "text", "text": para.text, "order": i} for i, para in enumerate(document.paragraphs) if para.text.strip()]
    for table in document.tables:
        rows = [[c.text.strip() for c in row.cells] for row in table.rows]
        blocks.append({"type": "table", "rows": rows, "headers": rows[0] if rows else [], "order": len(blocks)})
        for row in table.rows:
            cells = [c.text.strip() for c in row.cells]
            parts.append(" | ".join(cells))
    text = "\n".join(parts)
    page = NativePage(
        page_or_slide="1",
        raw_text=text,
        has_text_layer=bool(text.strip()),
        heading=path.stem,
        table_signals=bool(document.tables),
        language_hint=_language_hint(text),
        blocks=blocks,
        native_confidence=0.95 if text.strip() else 0.0,
        source_anchor=f"{path.name}, section 1",
        page_hash=hashlib.sha256((text + repr(blocks)).encode("utf-8")).hexdigest(),
    )
    _detect_formula_signals(page)
    _detect_questions(page)
    return [page]


def _parse_image(path: Path) -> list[NativePage]:
    # No native text for images: the visual path is required. No OCR by default.
    return [
        NativePage(
            page_or_slide="1",
            raw_text="",
            has_text_layer=False,
            has_images=True,
            heading=path.stem,
            source_anchor=f"{path.name}, image 1",
            page_hash=hashlib.sha256(path.read_bytes()).hexdigest(),
        )
    ]


def parse_native(path: Path, document_type: str, *, prefer_pymupdf: bool = True) -> list[NativePage]:
    """Native-first parsing. Raises on genuinely unreadable files (never fabricates)."""
    if document_type in {"txt", "md"}:
        return _parse_text_file(path)
    if document_type == "pdf":
        if prefer_pymupdf:
            try:
                return _parse_pdf_pymupdf(path)
            except Exception:
                return _parse_pdf(path)
        return _parse_pdf(path)
    if document_type == "pptx":
        return _parse_pptx(path)
    if document_type == "docx":
        return _parse_docx(path)
    if document_type in {"png", "jpg", "webp"}:
        return _parse_image(path)
    raise ValueError(f"no native parser for {document_type}")
