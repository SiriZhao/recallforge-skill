"""Fixture builders for Round 2 ingestion tests.

Each builder returns a real file on disk so the native parser and renderer run
against real bytes (no mocked parsing).
"""

from __future__ import annotations

import io
from pathlib import Path


def _insert_text(page, point, text, *, size=11):
    """Insert text with a CJK-capable built-in font so Chinese round-trips."""
    font = "china-s" if any("\u4e00" <= ch <= "\u9fff" for ch in text) else "helv"
    page.insert_text(point, text, fontsize=size, fontname=font)


def make_text_pdf(path: Path, *, title: str = "Titration Basics", body: str | None = None) -> Path:
    import pymupdf

    doc = pymupdf.open()
    page = doc.new_page()
    _insert_text(page, (72, 72), title, size=18)
    body = body or (
        "Standard solution is a solution of known concentration.\n"
        "c1 * V1 = c2 * V2\n"
        "Always write units and keep significant figures."
    )
    _insert_text(page, (72, 120), body)
    doc.save(str(path))
    doc.close()
    return path


def make_scanned_pdf(path: Path, *, text: str = "Scan OCR sample page") -> Path:
    """Create an image-only PDF (no text layer) - the 'scanned' case."""
    from PIL import Image, ImageDraw

    img = Image.new("RGB", (800, 1000), "white")
    draw = ImageDraw.Draw(img)
    draw.rectangle([60, 60, 740, 940], outline="black", width=2)
    draw.text((100, 120), text, fill="black")
    draw.text((100, 160), "handwritten style note here", fill="black")

    buf = io.BytesIO()
    img.save(buf, format="PNG")

    import pymupdf

    doc = pymupdf.open()
    page = doc.new_page(width=800, height=1000)
    page.insert_image(page.rect, stream=buf.getvalue())
    doc.save(str(path))
    doc.close()
    return path


def make_formula_pdf(path: Path, *, with_ambiguous: bool = True) -> Path:
    import pymupdf

    doc = pymupdf.open()
    page = doc.new_page()
    _insert_text(page, (72, 72), "Kinetics", size=18)
    lines = [
        "rate = k [A]^n [B]^m",
        "E_a / R = 2.303 * (T1 * T2) / (T2 - T1)",
        "H2SO4 + 2 NaOH -> Na2SO4 + 2 H2O",
    ]
    if with_ambiguous:
        lines.append("x_1 = (b +/- sqrt(b^2 - 4ac)) / 2a")
    _insert_text(page, (72, 130), "\n".join(lines))
    doc.save(str(path))
    doc.close()
    return path


def make_exam_pdf(path: Path) -> Path:
    import pymupdf

    doc = pymupdf.open()
    page = doc.new_page()
    content = (
        "Final Exam - Chemistry\n"
        "1. Which statement about standard solution is correct?\n"
        "A. concentration is not required to be accurate\n"
        "B. can be used for titration analysis\n"
        "C. does not need calibration\n"
        "D. must be used immediately\n"
        "(10 points)\n"
        "2. 简述系统误差和偶然误差的区别，并各举一例。(15 分)\n"
        "3. 用 0.1000 mol/L NaOH 滴定 25.00 mL HCl，消耗 24.80 mL，求 HCl 浓度。\n"
    )
    _insert_text(page, (72, 72), content)
    doc.save(str(path))
    doc.close()
    return path


def make_pptx(path: Path, *, slides: int = 2) -> Path:
    from pptx import Presentation
    from pptx.util import Inches

    prs = Presentation()
    prs.slides.add_slide(prs.slide_layouts[0])
    prs.slides[0].shapes.title.text = "Lecture: Standard Solution"
    prs.slides[0].placeholders[1].text = "Definition and calibration are key exam points."
    if slides > 1:
        slide2 = prs.slides.add_slide(prs.slide_layouts[1])
        slide2.shapes.title.text = "Titration procedure"
        body = slide2.placeholders[1].text_frame
        body.text = "Rinse, drain, add indicator, titrate, record volume."
    prs.save(str(path))
    return path


def make_docx(path: Path, *, with_table: bool = True) -> Path:
    import docx

    document = docx.Document()
    document.add_heading("Organic Chemistry Notes", level=1)
    document.add_paragraph("Functional groups and their reactions.")
    if with_table:
        table = document.add_table(rows=3, cols=2)
        table.rows[0].cells[0].text = "Functional group"
        table.rows[0].cells[1].text = "Reaction"
        table.rows[1].cells[0].text = "-OH"
        table.rows[1].cells[1].text = "esterification"
        table.rows[2].cells[0].text = "-COOH"
        table.rows[2].cells[1].text = "neutralization"
    document.save(str(path))
    return path


def make_image(path: Path, *, kind: str = "handwriting") -> Path:
    from PIL import Image, ImageDraw

    img = Image.new("RGB", (600, 400), "white")
    draw = ImageDraw.Draw(img)
    if kind == "handwriting":
        draw.text((40, 40), "C = n / V", fill="black")
        draw.text((40, 80), "remember sig figs", fill="black")
        draw.line([(40, 200), (300, 200)], fill="black", width=3)
    elif kind == "diagram":
        draw.ellipse([80, 80, 240, 240], outline="black", width=3)
        draw.line([(240, 160), (420, 160)], fill="black", width=2)
    img.save(str(path))
    return path


def make_mixed_language_pdf(path: Path) -> Path:
    import pymupdf

    doc = pymupdf.open()
    page = doc.new_page()
    _insert_text(page, (72, 72), "Bayes' theorem", size=16)
    _insert_text(page, (72, 120), "贝叶斯公式用于条件概率计算。")
    _insert_text(page, (72, 150), "P(A|B) = P(B|A) * P(A) / P(B)")
    doc.save(str(path))
    doc.close()
    return path
